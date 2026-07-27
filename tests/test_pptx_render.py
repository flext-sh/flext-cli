"""Model-driven PPTX rendering and reading contract tests."""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from flext_cli import cli, m
from flext_tests import tm


def test_pptx_render_empty_presentation() -> None:
    """Rendering an empty plan produces a valid PPTX."""
    plan = m.Cli.PptxPresentationPlan()
    result = cli.pptx_render(m.Cli.PptxRenderRequest(plan=plan))
    tm.that(result.success, eq=True, msg=result.error)
    tm.that(result.value.content[:2], eq=b"PK")
    presentation = Presentation(BytesIO(result.value.content))
    tm.that(len(presentation.slides), eq=0)


def test_pptx_render_with_slides() -> None:
    """Rendered slides contain the supplied titles."""
    plan = m.Cli.PptxPresentationPlan(
        slides=(m.Cli.PptxSlidePlan(title="Hello"), m.Cli.PptxSlidePlan(title="World"))
    )
    result = cli.pptx_render(m.Cli.PptxRenderRequest(plan=plan))
    tm.that(result.success, eq=True, msg=result.error)
    presentation = Presentation(BytesIO(result.value.content))
    tm.that(len(presentation.slides), eq=2)
    first_title = presentation.slides[0].shapes.title
    second_title = presentation.slides[1].shapes.title
    assert first_title is not None
    assert second_title is not None
    tm.that(first_title.text, eq="Hello")
    tm.that(second_title.text, eq="World")


def test_pptx_render_core_properties() -> None:
    """Core presentation properties are embedded in the rendered bytes."""
    plan = m.Cli.PptxPresentationPlan(
        slides=(m.Cli.PptxSlidePlan(title="Slide"),),
        core_properties={"title": "Test Title", "author": "Test Author"},
    )
    result = cli.pptx_render(m.Cli.PptxRenderRequest(plan=plan))
    tm.that(result.success, eq=True, msg=result.error)
    presentation = Presentation(BytesIO(result.value.content))
    tm.that(presentation.core_properties.title, eq="Test Title")
    tm.that(presentation.core_properties.author, eq="Test Author")


def test_pptx_read_round_trip() -> None:
    """Reading rendered bytes reproduces the presentation plan."""
    plan = m.Cli.PptxPresentationPlan(
        slides=(m.Cli.PptxSlidePlan(title="Hello"), m.Cli.PptxSlidePlan(title="World"))
    )
    render_result = cli.pptx_render(m.Cli.PptxRenderRequest(plan=plan))
    tm.that(render_result.success, eq=True, msg=render_result.error)
    read_result = cli.pptx_read(render_result.value.content)
    tm.that(read_result.success, eq=True, msg=read_result.error)
    tm.that(read_result.value.slides[0].title, eq="Hello")
    tm.that(read_result.value.slides[1].title, eq="World")


def test_pptx_read_rejects_invalid_bytes() -> None:
    """Reading invalid bytes returns a failure result."""
    result = cli.pptx_read(b"not a pptx")
    tm.that(result.success, eq=False)
    tm.that(result.error, ne=None)


def test_pptx_open_and_save() -> None:
    """Object-level byte boundary round-trips a presentation."""
    original = cli.Presentation()
    original.slides.add_slide(original.slide_layouts[0])
    source = BytesIO()
    original.save(source)
    open_result = cli.pptx_open(source.getvalue())
    tm.that(open_result.success, eq=True, msg=open_result.error)
    save_result = cli.pptx_save(open_result.value)
    tm.that(save_result.success, eq=True, msg=save_result.error)
    loaded = Presentation(BytesIO(save_result.value))
    tm.that(len(loaded.slides), eq=1)


def test_pptx_reexported_types() -> None:
    """python-pptx types are exposed through the generic boundary."""
    from pptx.dml.color import RGBColor as PptxRgbColor
    from pptx.enum.shapes import MSO_SHAPE as PPTX_MSO_SHAPE
    from pptx.util import Inches as PptxInches

    tm.that(cli.RGBColor, eq=PptxRgbColor)
    tm.that(cli.MSO_SHAPE, eq=PPTX_MSO_SHAPE)
    tm.that(cli.Inches, eq=PptxInches)
    tm.that(callable(cli.qn), eq=True)
